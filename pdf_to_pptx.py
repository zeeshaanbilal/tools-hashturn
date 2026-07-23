import sys
import fitz # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os

def main():
    if len(sys.argv) < 3:
        print("Usage: python pdf_to_pptx.py <input.pdf> <output.pptx>")
        sys.exit(1)
        
    input_pdf = sys.argv[1]
    output_pptx = sys.argv[2]
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    doc = fitz.open(input_pdf)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=150)
        img_path = f"temp_page_{page_num}_{os.path.basename(input_pdf)}.png"
        pix.save(img_path)
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width)
        
        os.remove(img_path)
        
    prs.save(output_pptx)
    doc.close()

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        import traceback
        print(f"Error: {e}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)
